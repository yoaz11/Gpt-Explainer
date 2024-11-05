from pptx import Presentation

def extract_text_from_pptx(file_path):
    """
    Extracts text from each slide in a PowerPoint (.pptx) file.

    Args:
        file_path (str): The path to the PowerPoint file.

    Returns:
        list: A list of strings, each representing the text from a slide.
    """
    prs = Presentation(file_path)
    slides_texts = []

    for slide in prs.slides:
        slide_text = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                slide_text += shape.text + " "
        slide_text = ' '.join(slide_text.split())  # Clean up whitespace
        if slide_text:  # Only add if slide has text
            slides_texts.append(slide_text)
    
    return slides_texts
